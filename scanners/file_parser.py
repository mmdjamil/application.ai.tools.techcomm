"""
file_parser.py

Extract text content line-by-line from:
- Word (.docx)
- PDF (.pdf)
- PowerPoint (.pptx)
- Excel (.xlsx)

Returns a list of dictionaries like:
[
    {"page": 1, "line": 1, "text": "some text"},
    {"page": 1, "line": 2, "text": "another line"},
]

Notes:
- PDF: page = actual PDF page number
- PPTX: page = slide number
- XLSX: page = sheet number
- DOCX: page = estimated page number (Word does not store reliable rendered page info)
"""

import io
import pdfplumber
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook


# -------------------------------------------------------------------
# Configuration
# -------------------------------------------------------------------
# For DOCX only: estimate page number by number of extracted lines.
# You can tune this value if needed.
ESTIMATED_LINES_PER_PAGE = 40


# -------------------------------------------------------------------
# DOCX parser
# -------------------------------------------------------------------
def parse_docx(file_bytes: bytes) -> list:
    """
    Parse a Word document (.docx).

    Since .docx does not provide reliable rendered page numbers through
    python-docx, this function estimates page number based on line count.

    Each non-empty paragraph = 1 line
    Each non-empty table cell = 1 line
    """
    doc = Document(io.BytesIO(file_bytes))
    results = []
    line_num = 0

    # Parse paragraphs
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            line_num += 1
            estimated_page = ((line_num - 1) // ESTIMATED_LINES_PER_PAGE) + 1
            results.append({
                "page": estimated_page,
                "line": line_num,
                "text": text
            })

    # Parse tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text = cell.text.strip()
                if text:
                    line_num += 1
                    estimated_page = ((line_num - 1) // ESTIMATED_LINES_PER_PAGE) + 1
                    results.append({
                        "page": estimated_page,
                        "line": line_num,
                        "text": text
                    })

    return results


# -------------------------------------------------------------------
# PDF parser
# -------------------------------------------------------------------
def parse_pdf(file_bytes: bytes) -> list:
    """
    Parse a PDF document page-by-page, line-by-line.
    Page number is accurate for PDFs.
    """
    results = []

    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page_num, page in enumerate(pdf.pages, start=1):
            text = page.extract_text()
            if text:
                for line_num, line in enumerate(text.split("\n"), start=1):
                    clean_line = line.strip()
                    if clean_line:
                        results.append({
                            "page": page_num,
                            "line": line_num,
                            "text": clean_line
                        })

    return results


# -------------------------------------------------------------------
# PPTX parser
# -------------------------------------------------------------------
def parse_pptx(file_bytes: bytes) -> list:
    """
    Parse a PowerPoint file (.pptx).
    Uses slide number as 'page'.
    """
    prs = Presentation(io.BytesIO(file_bytes))
    results = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        line_num = 0

        for shape in slide.shapes:
            # Text boxes and placeholders
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        line_num += 1
                        results.append({
                            "page": slide_num,
                            "line": line_num,
                            "text": text
                        })

            # Tables
            if hasattr(shape, "has_table") and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text:
                            line_num += 1
                            results.append({
                                "page": slide_num,
                                "line": line_num,
                                "text": text
                            })

    return results


# -------------------------------------------------------------------
# XLSX parser
# -------------------------------------------------------------------
def parse_xlsx(file_bytes: bytes) -> list:
    """
    Parse an Excel file (.xlsx).
    Uses sheet number as 'page'.
    Each non-empty row becomes one line.
    """
    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    results = []

    for sheet_num, sheet_name in enumerate(wb.sheetnames, start=1):
        ws = wb[sheet_name]

        for row_num, row in enumerate(ws.iter_rows(values_only=True), start=1):
            cell_texts = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            if cell_texts:
                combined_text = " | ".join(cell_texts)
                results.append({
                    "page": sheet_num,
                    "line": row_num,
                    "text": combined_text
                })

    wb.close()
    return results


# -------------------------------------------------------------------
# Main router
# -------------------------------------------------------------------
def parse_file(file_bytes: bytes, filename: str) -> list:
    """
    Route the uploaded file to the correct parser.
    Supported formats: .docx, .pdf, .pptx, .xlsx
    """
    ext = filename.lower().rsplit(".", 1)[-1]

    parsers = {
        "docx": parse_docx,
        "pdf": parse_pdf,
        "pptx": parse_pptx,
        "xlsx": parse_xlsx,
    }

    parser = parsers.get(ext)
    if parser is None:
        raise ValueError(
            f"Unsupported file type: .{ext}. Supported: .docx, .pdf, .pptx, .xlsx"
        )

    return parser(file_bytes)