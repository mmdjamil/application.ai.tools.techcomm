"""
document_rewriter.py

Re-open the original uploaded document and apply accepted (term, replacement)
edits per (page, line) key, returning the modified document in its original
binary format.

Supported formats:
  .docx  → python-docx  (preserves fonts, headings, tables, images)
  .pptx  → python-pptx  (preserves layout, theme, images)
  .xlsx  → openpyxl     (preserves sheets, styling, formulas in other cells)
  .pdf   → annotated PDF when possible, plain-text fallback otherwise

The walk order in each rewriter mirrors the corresponding parser in
file_parser.py exactly so that (page, line) keys align correctly.
"""

import io
import re
from copy import deepcopy

import fitz
from docx import Document
from docx.enum.text import WD_COLOR_INDEX
from docx.text.run import Run as DocxRun
from openpyxl import load_workbook
from openpyxl.comments import Comment
from openpyxl.styles import PatternFill
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.text.text import _Run as PptxRun

from .file_parser import ESTIMATED_LINES_PER_PAGE, parse_pdf
from .inclusive_scanner import (
    apply_all_accepted_replacements,
    pick_default_replacement,
)


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------

def _replacement_words_in_text(text: str, accepted: list) -> set[str]:
    """
    Return the set of default replacement words that are actually inserted into
    *text* when *accepted* replacements are applied.
    """
    replacement_words = set()
    for term, replacement in accepted:
        if re.search(r"\b" + re.escape(term) + r"\b", text, re.IGNORECASE):
            replacement_words.add(pick_default_replacement(replacement))
    return replacement_words


def _split_segments(new_text: str, replacement_words: set[str]) -> list[tuple[str, bool]]:
    """
    Split *new_text* into ``(text_chunk, is_highlight)`` segments, marking any
    whole-word occurrence of a replacement word for highlighting/styling.
    """
    if not new_text:
        return []
    if not replacement_words:
        return [(new_text, False)]

    pattern = re.compile(
        r"\b(?:"
        + "|".join(
            sorted(
                (re.escape(word) for word in replacement_words),
                key=len,
                reverse=True,
            )
        )
        + r")\b",
        re.IGNORECASE,
    )

    segments = []
    cursor = 0
    for match in pattern.finditer(new_text):
        if match.start() > cursor:
            segments.append((new_text[cursor:match.start()], False))
        segments.append((match.group(0), True))
        cursor = match.end()

    if cursor < len(new_text):
        segments.append((new_text[cursor:], False))

    return segments or [(new_text, False)]


def _ensure_docx_run(paragraph):
    if paragraph.runs:
        return paragraph.runs[0]
    return paragraph.add_run("")



def _apply_docx_runs(paragraph, segments: list[tuple[str, bool]]) -> None:
    """
    Write *segments* into a DOCX paragraph, highlighting only marked segments.
    """
    if not segments:
        return

    first_run = _ensure_docx_run(paragraph)
    original_runs = list(paragraph.runs)
    template_xml = deepcopy(first_run._r)
    previous_xml = first_run._r

    first_run.text = segments[0][0]
    first_run.font.highlight_color = (
        WD_COLOR_INDEX.YELLOW if segments[0][1] else None
    )

    for text, is_highlight in segments[1:]:
        new_xml = deepcopy(template_xml)
        previous_xml.addnext(new_xml)
        new_run = DocxRun(new_xml, paragraph)
        new_run.text = text
        new_run.font.highlight_color = WD_COLOR_INDEX.YELLOW if is_highlight else None
        previous_xml = new_xml

    for run in original_runs[1:]:
        run.text = ""
        run.font.highlight_color = None



def _ensure_pptx_run(paragraph):
    if paragraph.runs:
        return paragraph.runs[0]
    return paragraph.add_run()



def _style_pptx_run(run, is_highlight: bool) -> None:
    if not is_highlight:
        return
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xC0, 0x00, 0x00)



def _apply_pptx_runs(paragraph, segments: list[tuple[str, bool]]) -> None:
    """
    Write *segments* into a PPTX paragraph, styling only marked segments.
    """
    if not segments:
        return

    first_run = _ensure_pptx_run(paragraph)
    original_runs = list(paragraph.runs)
    template_xml = deepcopy(first_run._r)
    previous_xml = first_run._r

    first_run.text = segments[0][0]
    _style_pptx_run(first_run, segments[0][1])

    for text, is_highlight in segments[1:]:
        new_xml = deepcopy(template_xml)
        previous_xml.addnext(new_xml)
        new_run = PptxRun(new_xml, paragraph)
        new_run.text = text
        _style_pptx_run(new_run, is_highlight)
        previous_xml = new_xml

    for run in original_runs[1:]:
        run.text = ""



def _pdf_annotation_content(finding: dict, include_location_note: bool = False) -> str:
    lines = [
        f'Non-inclusive term: "{finding["found_word"]}"',
        (
            f'Suggested replacement: '
            f'"{pick_default_replacement(finding["suggested_replacement"])}"'
        ),
    ]
    if include_location_note:
        lines.append(f'(could not locate "{finding["found_word"]}" on this page)')
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Format-specific rewriters
# ---------------------------------------------------------------------------

def rewrite_docx(file_bytes: bytes, accepted_by_key: dict) -> bytes:
    """
    Re-open a .docx, apply accepted replacements per (page, line), return bytes.

    Walk order mirrors parse_docx: body paragraphs first, then table cells.
    Line numbering uses the same "non-empty after .strip()" rule as the parser
    to prevent (page, line) key drift.
    """
    doc = Document(io.BytesIO(file_bytes))
    line_num = 0

    # --- Body paragraphs (same order as parse_docx) ---
    for para in doc.paragraphs:
        if not para.text.strip():
            continue
        line_num += 1
        estimated_page = ((line_num - 1) // ESTIMATED_LINES_PER_PAGE) + 1
        key = (estimated_page, line_num)
        if key in accepted_by_key:
            run_text = "".join(run.text for run in para.runs)
            new_text = apply_all_accepted_replacements(run_text, accepted_by_key[key])
            if new_text != run_text:
                replacement_words = _replacement_words_in_text(
                    run_text, accepted_by_key[key]
                )
                _apply_docx_runs(para, _split_segments(new_text, replacement_words))

    # --- Table cells (same order as parse_docx) ---
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if not cell.text.strip():
                    continue
                line_num += 1
                estimated_page = ((line_num - 1) // ESTIMATED_LINES_PER_PAGE) + 1
                key = (estimated_page, line_num)
                if key in accepted_by_key:
                    # Apply the replacement to each paragraph in the cell
                    # individually; non-inclusive terms are single words and
                    # won't span cell-paragraph boundaries.
                    for cell_para in cell.paragraphs:
                        if not cell_para.text.strip():
                            continue
                        run_text = "".join(run.text for run in cell_para.runs)
                        new_text = apply_all_accepted_replacements(
                            run_text, accepted_by_key[key]
                        )
                        if new_text != run_text:
                            replacement_words = _replacement_words_in_text(
                                run_text, accepted_by_key[key]
                            )
                            _apply_docx_runs(
                                cell_para,
                                _split_segments(new_text, replacement_words),
                            )

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()



def rewrite_pptx(file_bytes: bytes, accepted_by_key: dict) -> bytes:
    """
    Re-open a .pptx, apply accepted replacements per (page, line), return bytes.

    Walk order mirrors parse_pptx: for each slide, text frames then tables.
    Line counter resets to 0 at the start of each slide (page = slide number).
    """
    prs = Presentation(io.BytesIO(file_bytes))

    for slide_num, slide in enumerate(prs.slides, start=1):
        line_num = 0

        for shape in slide.shapes:
            # --- Text frames (text boxes, placeholders) ---
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if not para.text.strip():
                        continue
                    line_num += 1
                    key = (slide_num, line_num)
                    if key in accepted_by_key:
                        run_text = "".join(run.text for run in para.runs)
                        new_text = apply_all_accepted_replacements(
                            run_text, accepted_by_key[key]
                        )
                        if new_text != run_text:
                            replacement_words = _replacement_words_in_text(
                                run_text, accepted_by_key[key]
                            )
                            _apply_pptx_runs(
                                para,
                                _split_segments(new_text, replacement_words),
                            )

            # --- Tables ---
            if hasattr(shape, "has_table") and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if not cell.text.strip():
                            continue
                        line_num += 1
                        key = (slide_num, line_num)
                        if key in accepted_by_key:
                            for cell_para in cell.text_frame.paragraphs:
                                if not cell_para.text.strip():
                                    continue
                                run_text = "".join(
                                    run.text for run in cell_para.runs
                                )
                                new_text = apply_all_accepted_replacements(
                                    run_text, accepted_by_key[key]
                                )
                                if new_text != run_text:
                                    replacement_words = _replacement_words_in_text(
                                        run_text, accepted_by_key[key]
                                    )
                                    _apply_pptx_runs(
                                        cell_para,
                                        _split_segments(new_text, replacement_words),
                                    )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()



def rewrite_xlsx(file_bytes: bytes, accepted_by_key: dict) -> bytes:
    """
    Re-open a .xlsx, apply accepted replacements per (page, line), return bytes.

    Opens in read-write / keep-formulas mode so untouched cells are preserved
    as-is (including formulas).  Walk order mirrors parse_xlsx: sheets in order,
    rows by enumerate index (row_num = actual 1-based position in the sheet,
    matching the parser's line numbering).
    """
    wb = load_workbook(io.BytesIO(file_bytes), read_only=False, data_only=False)

    for sheet_num, sheet_name in enumerate(wb.sheetnames, start=1):
        ws = wb[sheet_name]

        for row_num, row in enumerate(ws.iter_rows(), start=1):
            # Mirror the parser's non-empty row check
            has_content = any(
                cell.value is not None and str(cell.value).strip()
                for cell in row
            )
            if not has_content:
                continue

            key = (sheet_num, row_num)
            if key in accepted_by_key:
                for cell in row:
                    if cell.value is None or not isinstance(cell.value, str):
                        continue

                    original_value = cell.value
                    new_val = apply_all_accepted_replacements(
                        original_value, accepted_by_key[key]
                    )
                    if new_val != original_value:
                        comment_lines = []
                        for term, replacement in accepted_by_key[key]:
                            if re.search(
                                r"\b" + re.escape(term) + r"\b",
                                original_value,
                                re.IGNORECASE,
                            ):
                                comment_lines.append(
                                    f'Replaced "{term}" → '
                                    f'"{pick_default_replacement(replacement)}"'
                                )

                        cell.value = new_val
                        cell.fill = PatternFill(
                            start_color="FFFF00",
                            end_color="FFFF00",
                            fill_type="solid",
                        )
                        if comment_lines:
                            cell.comment = Comment("\n".join(comment_lines), "Doc Scanner")

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()



def _rewrite_pdf_as_text(file_bytes: bytes, accepted_by_key: dict) -> bytes:
    """
    Generate a plain-text fallback for PDFs.

    True in-place PDF rewriting (font embedding, glyph positioning) is out of
    scope.  Instead, the extracted text is re-assembled with replacements applied
    and returned as UTF-8 plain text, mirroring the original flat-text export.
    """
    parsed_lines = parse_pdf(file_bytes)
    chunks = []
    current_page = None

    for entry in parsed_lines:
        key = (entry["page"], entry["line"])
        text = entry["text"]
        if key in accepted_by_key:
            text = apply_all_accepted_replacements(text, accepted_by_key[key])

        if entry["page"] != current_page:
            if chunks:
                chunks.append("")
            chunks.append(f"--- Page/Section {entry['page']} ---")
            current_page = entry["page"]

        chunks.append(f"L{entry['line']:>4}: {text}")

    result = "\n".join(chunks).strip() + "\n" if chunks else ""
    return result.encode("utf-8")



def rewrite_pdf_with_annotations(
    file_bytes: bytes, accepted_findings: list[dict]
) -> bytes:
    """
    Add highlight + text annotations for each accepted PDF finding.
    """
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    try:
        for finding in accepted_findings:
            page_number = int(finding["page"]) - 1
            if page_number < 0 or page_number >= len(doc):
                continue

            page = doc[page_number]
            rects = page.search_for(finding["found_word"])
            content = _pdf_annotation_content(finding)

            if rects:
                for rect in rects:
                    highlight = page.add_highlight_annot(rect)
                    highlight.set_info(title="Doc Scanner")
                    highlight.update()

                    note = page.add_text_annot(rect.tl, content)
                    note.set_info(title="Doc Scanner")
                    note.update()
            else:
                note = page.add_text_annot(
                    fitz.Point(36, 36),
                    _pdf_annotation_content(finding, include_location_note=True),
                )
                note.set_info(title="Doc Scanner")
                note.update()

        return doc.tobytes()
    finally:
        doc.close()


# ---------------------------------------------------------------------------
# Public router
# ---------------------------------------------------------------------------

def rewrite_file(
    file_bytes: bytes,
    filename: str,
    accepted_by_key: dict,
    accepted_findings: list[dict] | None = None,
) -> tuple:
    """
    Apply accepted (term, replacement) edits to *file_bytes* and return
    ``(output_bytes, output_filename, mime_type)``.

    For .pdf the function returns an annotated PDF when *accepted_findings* are
    provided and annotation succeeds. Otherwise it falls back to a plain-text
    export for compatibility with existing callers.
    """
    base = filename.rsplit(".", 1)[0]
    ext = filename.lower().rsplit(".", 1)[-1]

    if ext == "docx":
        return (
            rewrite_docx(file_bytes, accepted_by_key),
            f"{base}_corrected.docx",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )
    if ext == "pptx":
        return (
            rewrite_pptx(file_bytes, accepted_by_key),
            f"{base}_corrected.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    if ext == "xlsx":
        return (
            rewrite_xlsx(file_bytes, accepted_by_key),
            f"{base}_corrected.xlsx",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )
    if ext == "pdf":
        if accepted_findings is not None:
            try:
                return (
                    rewrite_pdf_with_annotations(file_bytes, accepted_findings),
                    f"{base}_corrected.pdf",
                    "application/pdf",
                )
            except Exception:
                pass
        return (
            _rewrite_pdf_as_text(file_bytes, accepted_by_key),
            f"{base}_corrected.txt",
            "text/plain",
        )

    raise ValueError(
        f"Unsupported file type: .{ext}. Supported: .docx, .pptx, .xlsx, .pdf"
    )
