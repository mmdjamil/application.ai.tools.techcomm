"""
test_document_rewriter.py

Round-trip tests for scanners.document_rewriter.  Each test builds a minimal
in-memory document, runs rewrite_file, re-opens the result, and asserts that
the accepted replacements were applied while the document remains well-formed.
"""

import io
import unittest
from unittest.mock import patch

import fitz
from docx import Document as DocxDocument
from docx.enum.text import WD_COLOR_INDEX
from pptx import Presentation
from pptx.util import Inches
from openpyxl import load_workbook, Workbook

from scanners.document_rewriter import rewrite_file
from scanners.file_parser import parse_docx


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_docx(*paragraphs: str) -> bytes:
    """Return bytes of a .docx containing the given paragraphs."""
    doc = DocxDocument()
    for text in paragraphs:
        doc.add_paragraph(text)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_pptx(text: str) -> bytes:
    """Return bytes of a .pptx with one slide and one text box."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # blank layout
    slide = prs.slides.add_slide(slide_layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    txBox.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_xlsx(value: str) -> bytes:
    """Return bytes of a .xlsx with one sheet and one cell containing *value*."""
    wb = Workbook()
    ws = wb.active
    ws["A1"] = value
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------

class DocumentRewriterTests(unittest.TestCase):

    # --- DOCX ---

    def test_docx_replaces_blacklist_and_whitelist(self):
        """Round-trip: blacklist → denylist, Whitelist → Allowlist."""
        file_bytes = _make_docx("The blacklist is bad.", "Whitelist all trusted IPs.")
        # Parse to get (page, line) keys
        parsed = parse_docx(file_bytes)
        # Both lines should be found
        self.assertEqual(len(parsed), 2)
        key1 = (parsed[0]["page"], parsed[0]["line"])
        key2 = (parsed[1]["page"], parsed[1]["line"])

        accepted_by_key = {
            key1: [("blacklist", "denylist")],
            key2: [("whitelist", "allowlist")],
        }

        output_bytes, output_filename, mime = rewrite_file(
            file_bytes, "report.docx", accepted_by_key
        )

        self.assertTrue(output_filename.endswith(".docx"))
        self.assertEqual(
            mime,
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )

        # Re-open and verify text
        result_doc = DocxDocument(io.BytesIO(output_bytes))
        texts = [p.text for p in result_doc.paragraphs if p.text.strip()]
        self.assertIn("denylist", texts[0])
        self.assertNotIn("blacklist", texts[0])
        self.assertIn("Allowlist", texts[1])
        self.assertNotIn("Whitelist", texts[1])

    def test_docx_document_opens_without_error(self):
        """The rewritten .docx must be a valid document that opens cleanly."""
        file_bytes = _make_docx("Use the blacklist feature carefully.")
        parsed = parse_docx(file_bytes)
        key = (parsed[0]["page"], parsed[0]["line"])
        accepted_by_key = {key: [("blacklist", "denylist")]}

        output_bytes, _, _ = rewrite_file(file_bytes, "doc.docx", accepted_by_key)

        # Should not raise
        result_doc = DocxDocument(io.BytesIO(output_bytes))
        self.assertTrue(len(result_doc.paragraphs) > 0)

    def test_docx_highlights_replaced_word(self):
        """Replaced DOCX words should be highlighted in yellow."""
        file_bytes = _make_docx("The blacklist is bad.")
        parsed = parse_docx(file_bytes)
        key = (parsed[0]["page"], parsed[0]["line"])
        accepted_by_key = {key: [("blacklist", "denylist")]}

        output_bytes, _, _ = rewrite_file(file_bytes, "report.docx", accepted_by_key)

        result_doc = DocxDocument(io.BytesIO(output_bytes))
        highlighted_runs = [
            run
            for para in result_doc.paragraphs
            for run in para.runs
            if run.font.highlight_color == WD_COLOR_INDEX.YELLOW
            and "denylist" in run.text.lower()
        ]
        self.assertTrue(highlighted_runs)

    # --- PPTX ---

    def test_pptx_replaces_master(self):
        """Round-trip: master → primary in a single-slide presentation."""
        file_bytes = _make_pptx("The master branch is here.")
        accepted_by_key = {(1, 1): [("master", "primary / initiator")]}

        output_bytes, output_filename, mime = rewrite_file(
            file_bytes, "slides.pptx", accepted_by_key
        )

        self.assertTrue(output_filename.endswith(".pptx"))
        self.assertEqual(
            mime,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

        prs = Presentation(io.BytesIO(output_bytes))
        slide = prs.slides[0]
        all_text = " ".join(
            para.text
            for shape in slide.shapes
            if hasattr(shape, "has_text_frame") and shape.has_text_frame
            for para in shape.text_frame.paragraphs
        )
        self.assertIn("primary", all_text)
        self.assertNotIn("master", all_text)

    def test_pptx_styles_replaced_word(self):
        """Replaced PPTX words should be visibly styled."""
        file_bytes = _make_pptx("The master branch is here.")
        accepted_by_key = {(1, 1): [("master", "primary / initiator")]}

        output_bytes, _, _ = rewrite_file(file_bytes, "slides.pptx", accepted_by_key)

        prs = Presentation(io.BytesIO(output_bytes))
        styled_found = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if not (hasattr(shape, "has_text_frame") and shape.has_text_frame):
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "primary" not in run.text.lower():
                            continue
                        has_highlight_xml = "<a:highlight" in run._r.xml
                        has_bold = run.font.bold is True
                        has_color = run.font.color.rgb is not None
                        if has_highlight_xml or has_bold or has_color:
                            styled_found = True
        self.assertTrue(styled_found)

    # --- XLSX ---

    def test_xlsx_replaces_slave(self):
        """Round-trip: slave → secondary in a single-cell workbook."""
        file_bytes = _make_xlsx("slave")
        # The parser uses row_num=1 for the first row, sheet_num=1
        accepted_by_key = {(1, 1): [("slave", "secondary / target")]}

        output_bytes, output_filename, mime = rewrite_file(
            file_bytes, "data.xlsx", accepted_by_key
        )

        self.assertTrue(output_filename.endswith(".xlsx"))
        self.assertEqual(
            mime,
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )

        wb = load_workbook(io.BytesIO(output_bytes))
        ws = wb.active
        self.assertEqual(ws["A1"].value, "secondary")

    def test_xlsx_highlights_and_comments_cell(self):
        """Changed XLSX cells should be filled yellow and receive a comment."""
        file_bytes = _make_xlsx("blacklist master")
        accepted_by_key = {
            (1, 1): [
                ("blacklist", "denylist"),
                ("master", "primary / initiator"),
            ]
        }

        output_bytes, _, _ = rewrite_file(file_bytes, "data.xlsx", accepted_by_key)

        wb = load_workbook(io.BytesIO(output_bytes))
        ws = wb.active
        self.assertTrue(ws["A1"].fill.start_color.rgb.endswith("FFFF00"))
        self.assertIsNotNone(ws["A1"].comment)
        self.assertIn("blacklist", ws["A1"].comment.text)
        self.assertIn("denylist", ws["A1"].comment.text)
        self.assertIn("master", ws["A1"].comment.text)
        self.assertIn("primary", ws["A1"].comment.text)

    # --- PDF fallback ---

    def test_pdf_returns_pdf_with_annotations(self):
        """Accepted PDF findings should produce an annotated PDF."""
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((72, 72), "The blacklist is bad.")
        file_bytes = doc.tobytes()
        doc.close()

        accepted_by_key = {(1, 1): [("blacklist", "denylist")]}
        accepted_findings = [
            {
                "page": 1,
                "line": 1,
                "found_word": "blacklist",
                "term": "blacklist",
                "suggested_replacement": "denylist",
            }
        ]

        output_bytes, output_filename, mime = rewrite_file(
            file_bytes,
            "report.pdf",
            accepted_by_key,
            accepted_findings=accepted_findings,
        )

        self.assertEqual(mime, "application/pdf")
        self.assertTrue(output_filename.endswith(".pdf"))

        result = fitz.open(stream=output_bytes, filetype="pdf")
        page = result[0]
        annot_types = []
        annot_contents = []
        for annot in page.annots():
            annot_types.append(annot.type[1])
            annot_contents.append(annot.info.get("content") or "")
        highlight_count = sum(annot_type == "Highlight" for annot_type in annot_types)
        text_count = sum(annot_type == "Text" for annot_type in annot_types)
        self.assertGreaterEqual(highlight_count, 1)
        self.assertGreaterEqual(text_count, 1)
        self.assertTrue(any("denylist" in content for content in annot_contents))
        result.close()

    def test_pdf_fallback_returns_text_plain(self):
        """PDF path must return mime=text/plain and apply the replacement."""
        mock_lines = [
            {"page": 1, "line": 1, "text": "The slave device is connected."},
        ]
        accepted_by_key = {(1, 1): [("slave", "secondary / target")]}

        with patch("scanners.document_rewriter.parse_pdf", return_value=mock_lines):
            output_bytes, output_filename, mime = rewrite_file(
                b"fake-pdf-bytes", "report.pdf", accepted_by_key
            )

        self.assertEqual(mime, "text/plain")
        self.assertTrue(output_filename.endswith(".txt"))
        decoded = output_bytes.decode("utf-8")
        self.assertIn("secondary", decoded)
        self.assertNotIn("slave", decoded)

    # --- Line drift guard ---

    def test_docx_line_drift_empty_paragraph_between_targets(self):
        """
        An empty paragraph between two target paragraphs must not shift
        (page, line) keys — the rewriter and parser must agree on numbering.
        """
        doc = DocxDocument()
        doc.add_paragraph("The blacklist is here.")  # line 1
        doc.add_paragraph("")                         # empty — skipped by both
        doc.add_paragraph("The whitelist is there.")  # line 2

        buf = io.BytesIO()
        doc.save(buf)
        file_bytes = buf.getvalue()

        # Verify the parser assigns line 1 and 2 (not 1 and 3)
        parsed = parse_docx(file_bytes)
        self.assertEqual(len(parsed), 2)
        self.assertEqual(parsed[0]["line"], 1)
        self.assertEqual(parsed[1]["line"], 2)

        key1 = (parsed[0]["page"], parsed[0]["line"])
        key2 = (parsed[1]["page"], parsed[1]["line"])

        accepted_by_key = {
            key1: [("blacklist", "denylist")],
            key2: [("whitelist", "allowlist")],
        }

        output_bytes, _, _ = rewrite_file(file_bytes, "test.docx", accepted_by_key)

        result_doc = DocxDocument(io.BytesIO(output_bytes))
        result_texts = [p.text for p in result_doc.paragraphs if p.text.strip()]

        self.assertIn("denylist", result_texts[0])
        self.assertNotIn("blacklist", result_texts[0])
        self.assertIn("allowlist", result_texts[1])
        self.assertNotIn("whitelist", result_texts[1])


if __name__ == "__main__":
    unittest.main()
